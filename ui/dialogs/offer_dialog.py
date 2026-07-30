import os
import io
import pandas as pd
from PIL import Image as PILImage
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import mm
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER
from reportlab.pdfbase.pdfmetrics import stringWidth
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Image as RLImage, Spacer

# pyrefly: ignore [missing-import]
import arabic_reshaper
from bidi.algorithm import get_display

def shape_arabic(text):
    if not isinstance(text, str):
        text = str(text)
    if not text:
        return ""
    return get_display(arabic_reshaper.reshape(text))

from PyQt6.QtWidgets import (QDialog, QVBoxLayout, QHBoxLayout, QLabel, QPushButton,
                             QHeaderView, QTableWidgetItem, QMessageBox, QFileDialog, QWidget, QInputDialog, QLineEdit, QCheckBox, QAbstractItemView)
from PyQt6.QtCore import Qt, QEvent
from PyQt6.QtGui import QColor, QFont

from ui.widgets.dynamic_table import DynamicTableWidget
from ui.dialogs.mapping_dialog import MappingDialog
from ui.dialogs.pdf_dialog import PdfSettingsDialog
from ui.dialogs.set_selection_dialog import SetSelectionDialog
from core.utils import show_loading

class OfferListDialog(QDialog):
    def __init__(self, master_tab, offer_data, parent=None):
        super().__init__(parent)
        self.master_tab = master_tab
        self.offer_data = offer_data
        self.extra_columns = []
        self.setWindowTitle("Manage Offer List")
        self.resize(800, 600)
        self.init_ui()
        self.refresh_offer_table()
        self.offer_table.scrollToBottom()

    def init_ui(self):
        layout = QVBoxLayout(self)
        layout.setContentsMargins(12, 12, 12, 12)
        layout.setSpacing(12)

        layout.addWidget(QLabel("<b>Current Offer List</b>"))
        
        # Search Layout
        search_layout = QHBoxLayout()
        self.search_input = QLineEdit()
        self.search_input.setPlaceholderText("Search in offer list (Press Enter for next)...")
        self.search_input.returnPressed.connect(self.find_next)
        self.search_input.textChanged.connect(self.reset_search)
        
        self.btn_find_next = QPushButton("Find Next")
        self.btn_find_next.clicked.connect(self.find_next)
        
        self.chk_headers_only = QCheckBox("In Headers")
        self.chk_headers_only.stateChanged.connect(self.reset_search)
        
        search_layout.addWidget(self.search_input)
        search_layout.addWidget(self.btn_find_next)
        search_layout.addWidget(self.chk_headers_only)
        layout.addLayout(search_layout)

        self.current_search_row = -1

        self.offer_table = DynamicTableWidget()
        self.offer_table.setColumnCount(4)
        self.offer_table.setHorizontalHeaderLabels(["", "CODE", "DESCRIPTION", "QTY"])
        self.offer_table.setColumnWidth(0, 80)
        self.offer_table.horizontalHeader().setSectionResizeMode(1, QHeaderView.ResizeMode.ResizeToContents)
        self.offer_table.horizontalHeader().setSectionResizeMode(2, QHeaderView.ResizeMode.Stretch)
        self.offer_table.horizontalHeader().setSectionResizeMode(3, QHeaderView.ResizeMode.ResizeToContents)
        self.offer_table.horizontalHeader().setSectionsMovable(True)
        self.offer_table.setStyleSheet("""
            QTableWidget::item:selected {
                background-color: #FBBF24;
                color: #000000;
            }
        """)
        self.offer_table.itemChanged.connect(self.on_offer_item_changed)
        self.offer_table.currentCellChanged.connect(self._on_table_selection_changed)
        self.offer_table.installEventFilter(self)
        layout.addWidget(self.offer_table)

        btn_layout = QHBoxLayout()
        self.btn_clear = QPushButton("Clear List")
        self.btn_clear.clicked.connect(self.clear_offer_list)
        
        self.btn_add_section = QPushButton("Add Section Header")
        self.btn_add_section.clicked.connect(self.add_section_manually)
        
        self.btn_bulk_upload = QPushButton("Bulk Upload (Excel/CSV)")
        self.btn_bulk_upload.clicked.connect(self.bulk_upload_offer)

        self.btn_load_set = QPushButton("Load Set")
        self.btn_load_set.clicked.connect(self.load_set_offer)

        self.btn_export_excel = QPushButton("Export to Excel")
        self.btn_export_excel.clicked.connect(self.export_excel)
        
        self.btn_export_pdf = QPushButton("Export to PDF")
        self.btn_export_pdf.clicked.connect(self.export_pdf)
        
        self.btn_export_pptx = QPushButton("Export to PPTX")
        self.btn_export_pptx.setObjectName("primaryButton")
        self.btn_export_pptx.clicked.connect(self.export_pptx)

        btn_layout.addWidget(self.btn_clear)
        btn_layout.addWidget(self.btn_add_section)
        btn_layout.addWidget(self.btn_bulk_upload)
        btn_layout.addWidget(self.btn_load_set)
        btn_layout.addStretch()
        btn_layout.addWidget(self.btn_export_excel)
        btn_layout.addWidget(self.btn_export_pdf)
        btn_layout.addWidget(self.btn_export_pptx)
        layout.addLayout(btn_layout)

    def refresh_offer_table(self):
        self.offer_table.blockSignals(True)
        self.offer_table.clearSpans()
        self.offer_table.setColumnCount(5 + len(self.extra_columns))
        self.offer_table.setHorizontalHeaderLabels(["", "#", "CODE", "DESCRIPTION", "QTY"] + [c.upper() for c in self.extra_columns])
        self.offer_table.setRowCount(len(self.offer_data))
        
        serial_counter = 1
        for i, data in enumerate(self.offer_data):
            btn_up = QPushButton("▲")
            btn_up.setFixedSize(20, 20)
            btn_up.setCursor(Qt.CursorShape.PointingHandCursor)
            btn_up.clicked.connect(lambda _, row=i: self.move_offer_item_up(row))
            if i == 0: btn_up.setEnabled(False)

            btn_down = QPushButton("▼")
            btn_down.setFixedSize(20, 20)
            btn_down.setCursor(Qt.CursorShape.PointingHandCursor)
            btn_down.clicked.connect(lambda _, row=i: self.move_offer_item_down(row))
            if i == len(self.offer_data) - 1: btn_down.setEnabled(False)

            btn = QPushButton("−")
            btn.setObjectName("removeButton")
            btn.setFixedSize(20, 20)
            btn.setCursor(Qt.CursorShape.PointingHandCursor)
            btn.setToolTip("Remove from offer list")
            btn.clicked.connect(lambda _, row=i: self.remove_offer_item(row))

            btn_container = QWidget()
            btn_container.setStyleSheet("QWidget { background: transparent; }")
            btn_layout = QHBoxLayout(btn_container)
            btn_layout.setContentsMargins(2, 0, 2, 0)
            btn_layout.setSpacing(2)
            btn_layout.addWidget(btn_up)
            btn_layout.addWidget(btn_down)
            btn_layout.addWidget(btn)
            self.offer_table.setRowHeight(i, 36)
            self.offer_table.setCellWidget(i, 0, btn_container)

            if data.get("is_section", False):
                serial_counter = 1
                for c in range(1, 5 + len(self.extra_columns)):
                    if c == 3:
                        item = QTableWidgetItem(data["desc"])
                    else:
                        item = QTableWidgetItem("")
                    item.setFlags(Qt.ItemFlag.ItemIsSelectable | Qt.ItemFlag.ItemIsEditable | Qt.ItemFlag.ItemIsEnabled)
                    item.setFont(QFont("Segoe UI", 10, QFont.Weight.Bold))
                    item.setBackground(QColor("#DDE3EC"))
                    if c == 3:
                        item.setTextAlignment(Qt.AlignmentFlag.AlignCenter)
                    self.offer_table.setItem(i, c, item)
            else:
                item_sr = QTableWidgetItem(str(serial_counter))
                item_sr.setFlags(Qt.ItemFlag.ItemIsSelectable | Qt.ItemFlag.ItemIsEditable | Qt.ItemFlag.ItemIsEnabled)
                item_sr.setTextAlignment(Qt.AlignmentFlag.AlignCenter)
                self.offer_table.setItem(i, 1, item_sr)
                
                item_code = QTableWidgetItem(data["code"])
                item_code.setFlags(Qt.ItemFlag.ItemIsSelectable | Qt.ItemFlag.ItemIsEditable | Qt.ItemFlag.ItemIsEnabled)
                self.offer_table.setItem(i, 2, item_code)

                item_desc = QTableWidgetItem(data["desc"])
                item_desc.setFlags(Qt.ItemFlag.ItemIsSelectable | Qt.ItemFlag.ItemIsEditable | Qt.ItemFlag.ItemIsEnabled)
                self.offer_table.setItem(i, 3, item_desc)

                item_qty = QTableWidgetItem(str(data["qty"]))
                item_qty.setFlags(Qt.ItemFlag.ItemIsSelectable | Qt.ItemFlag.ItemIsEditable | Qt.ItemFlag.ItemIsEnabled)
                self.offer_table.setItem(i, 4, item_qty)

                for idx, col_name in enumerate(self.extra_columns):
                    item_extra = QTableWidgetItem(str(data.get(col_name, "")))
                    item_extra.setFlags(Qt.ItemFlag.ItemIsSelectable | Qt.ItemFlag.ItemIsEditable | Qt.ItemFlag.ItemIsEnabled)
                    self.offer_table.setItem(i, 5 + idx, item_extra)
                    
                serial_counter += 1

        self.offer_table.blockSignals(False)

    def reset_search(self):
        self.current_search_row = -1
        if not self.search_input.text().strip():
            self.offer_table.clearSelection()

    def find_next(self):
        search_text = self.search_input.text().strip().lower()
        if not search_text or not self.offer_data:
            return
            
        headers_only = self.chk_headers_only.isChecked()
        start_row = self.current_search_row + 1
        
        # Loop around the table to find the next match
        for offset in range(len(self.offer_data)):
            row = (start_row + offset) % len(self.offer_data)
            data = self.offer_data[row]
            is_section = data.get("is_section", False)
            
            if headers_only and not is_section:
                continue
                
            match = False
            if headers_only:
                if search_text in str(data.get("desc", "")).lower():
                    match = True
            else:
                for val in data.values():
                    if search_text in str(val).lower():
                        match = True
                        break
                        
            if match:
                self.current_search_row = row
                self.offer_table.selectRow(row)
                # Scroll to the selected row, specifically column 1 (to ensure visibility)
                self.offer_table.scrollToItem(self.offer_table.item(row, 1), QAbstractItemView.ScrollHint.PositionAtCenter)
                return
                
    def eventFilter(self, source, event):
        if source == self.offer_table and event.type() == QEvent.Type.KeyPress:
            if event.key() == Qt.Key.Key_Return or event.key() == Qt.Key.Key_Enter:
                selected_rows = sorted(list(set(item.row() for item in self.offer_table.selectedItems())))
                if not selected_rows:
                    current_row = self.offer_table.currentRow()
                    if current_row >= 0:
                        selected_rows = [current_row]
                
                if selected_rows:
                    import copy
                    for row in selected_rows:
                        if 0 <= row < len(self.offer_data):
                            item_copy = copy.deepcopy(self.offer_data[row])
                            self.offer_data.append(item_copy)
                    
                    if hasattr(self, 'master_tab') and hasattr(self.master_tab, 'save_offer_list'):
                        self.master_tab.save_offer_list()
                    self.refresh_offer_table()
                    return True
        return super().eventFilter(source, event)

    def on_offer_item_changed(self, item):
        row = item.row()
        col = item.column()

        if row < 0 or row >= len(self.offer_data):
            return

        new_val = item.text().strip()

        if col == 2:
            self.offer_data[row]["code"] = new_val
        elif col == 3:
            self.offer_data[row]["desc"] = new_val
        elif col == 4:
            try:
                self.offer_data[row]["qty"] = int(new_val)
            except ValueError:
                pass
        elif col >= 5:
            col_name = self.extra_columns[col - 5]
            self.offer_data[row][col_name] = new_val
        
        if hasattr(self, 'master_tab') and hasattr(self.master_tab, 'save_offer_list'):
            self.master_tab.save_offer_list()

    def _on_table_selection_changed(self, current_row, current_col, prev_row, prev_col):
        if current_row < 0 or current_row >= len(self.offer_data):
            return
        data = self.offer_data[current_row]
        if data.get("is_section", False):
            return
        code = data.get("code", "")
        
        master_df = self.master_tab.df
        if master_df is None or master_df.empty:
            return
            
        mask = master_df['code'].astype(str).str.replace('-', '').str.strip().str.lower() == str(code).replace('-', '').strip().lower()
        match = master_df[mask]
        if not match.empty:
            main_win = self.master_tab.window()
            if hasattr(main_win, 'side_panel'):
                main_win.side_panel.update_panel(code)

    def move_offer_item_up(self, row_idx):
        if 0 < row_idx < len(self.offer_data):
            self.offer_data[row_idx - 1], self.offer_data[row_idx] = self.offer_data[row_idx], self.offer_data[row_idx - 1]
            if hasattr(self, 'master_tab') and hasattr(self.master_tab, 'save_offer_list'):
                self.master_tab.save_offer_list()
            self.refresh_offer_table()

    def move_offer_item_down(self, row_idx):
        if 0 <= row_idx < len(self.offer_data) - 1:
            self.offer_data[row_idx + 1], self.offer_data[row_idx] = self.offer_data[row_idx], self.offer_data[row_idx + 1]
            if hasattr(self, 'master_tab') and hasattr(self.master_tab, 'save_offer_list'):
                self.master_tab.save_offer_list()
            self.refresh_offer_table()

    def remove_offer_item(self, row_idx):
        if 0 <= row_idx < len(self.offer_data):
            self.offer_data.pop(row_idx)
            if hasattr(self, 'master_tab') and hasattr(self.master_tab, 'save_offer_list'):
                self.master_tab.save_offer_list()
            self.refresh_offer_table()

    def clear_offer_list(self):
        msg_box = QMessageBox(self)
        msg_box.setWindowTitle("Clear Offer List")
        msg_box.setText("Choose clearing options:")
        btn_clear_only = msg_box.addButton("Clear Offer List Only", QMessageBox.ButtonRole.ActionRole)
        btn_clear_both = msg_box.addButton("Clear Offer List & Reset Mapping", QMessageBox.ButtonRole.ActionRole)
        btn_cancel = msg_box.addButton(QMessageBox.StandardButton.Cancel)
        msg_box.exec()
        
        if msg_box.clickedButton() == btn_cancel:
            return

        self.offer_data.clear()
        if hasattr(self, 'master_tab') and hasattr(self.master_tab, 'save_offer_list'):
            self.master_tab.save_offer_list()
        self.refresh_offer_table()
        
        if msg_box.clickedButton() == btn_clear_both:
            if hasattr(self.master_tab, 'wizard_panel'):
                self.master_tab.wizard_panel.reset_mapping()

    def add_section_manually(self):
        text, ok = QInputDialog.getText(self, "Add Section Header", "Enter Header Title:")
        if ok and text.strip():
            self.offer_data.append({"is_section": True, "desc": text.strip()})
            if hasattr(self, 'master_tab') and hasattr(self.master_tab, 'save_offer_list'):
                self.master_tab.save_offer_list()
            self.refresh_offer_table()

    def load_set_offer(self):
        master_df = self.master_tab.df
        if master_df is None or master_df.empty:
            QMessageBox.warning(self, "No Master Data", "Please load the master catalogue before loading sets.")
            return
            
        dlg = SetSelectionDialog(self.master_tab, parent=self)
        if dlg.exec() != QDialog.DialogCode.Accepted:
            return
            
        set_items = dlg.get_selected_items()
        if set_items is None or set_items.empty:
            QMessageBox.information(self, "No Items", "The selected set has no items.")
            return

        with show_loading(self, "Merging Set Data..."):
            master_copy = master_df.copy()
            master_copy['join_key'] = master_copy['code'].astype(str).str.replace('-', '').str.strip().str.lower()
            
            set_items_copy = set_items.copy()
            set_items_copy['join_key'] = set_items_copy['Item_SKU'].astype(str).str.replace('-', '').str.strip().str.lower()
            
            # Relational JOIN operation using Pandas
            merged_df = pd.merge(set_items_copy, master_copy, on='join_key', how='left')
            
            set_name_series = dlg.unique_sets_df[dlg.unique_sets_df['Set_Code'].astype(str) == dlg.selected_set_code]['Set_Name']
            set_name = set_name_series.iloc[0] if not set_name_series.empty else dlg.selected_set_code
            self.offer_data.append({"is_section": True, "desc": f"Set: {set_name} ({dlg.selected_set_code})"})
            
            skipped = []
            for _, row in merged_df.iterrows():
                code = str(row.get('Item_SKU', ''))
                if not code or code.lower() == 'nan':
                    continue
                    
                qty = row.get('Quantity', 1)
                try:
                    qty = int(float(qty))
                except (ValueError, TypeError):
                    qty = 1
                
                desc = row.get('description', None)
                if pd.isna(desc) or not str(desc).strip():
                    desc = str(row.get('Item_Name', ''))
                    skipped.append(code)
                else:
                    desc = str(desc)
                    
                self.offer_data.append({
                    "code": code,
                    "desc": desc,
                    "qty": qty
                })
        
        self.refresh_offer_table()
        if hasattr(self, 'master_tab') and hasattr(self.master_tab, 'save_offer_list'):
            self.master_tab.save_offer_list()
        if skipped:
            QMessageBox.warning(
                self,
                "Missing Master Data",
                f"{len(skipped)} items from the set were not found in the master catalogue. They have been added using their default names from the sets export."
            )

    def bulk_upload_offer(self):
        file_path, _ = QFileDialog.getOpenFileName(
            self,
            "Open Offer Upload File",
            "",
            "Data Files (*.csv *.xlsx *.xls)"
        )
        if not file_path:
            return
        dlg = MappingDialog(file_path, required_fields=['code'], optional_fields=['description', 'quantity'], allow_extras=True, parent=self)
        if dlg.exec() != QDialog.DialogCode.Accepted:
            return

        with show_loading(self, "Processing Bulk Upload..."):
            try:
                if file_path.lower().endswith('.csv'):
                    raw_df = pd.read_csv(file_path, header=dlg.header_row)
                else:
                    raw_df = pd.read_excel(file_path, sheet_name=dlg.selected_sheet, header=dlg.header_row)

                rename_map = {'code': 'code', 'description': 'desc', 'quantity': 'qty'}
                mapped_cols = {}
                for source_col, target_key in rename_map.items():
                    if source_col in dlg.mappings and dlg.mappings[source_col]:
                        mapped_cols[source_col] = dlg.mappings[source_col]
                    else:
                        mapped_cols[source_col] = None

                if mapped_cols['code'] is None:
                    QMessageBox.warning(self, "Invalid Mapping", "Please map at least the 'code' column.")
                    return

                new_extras = getattr(dlg, 'extras', [])
                for ext in new_extras:
                    if ext not in self.extra_columns:
                        self.extra_columns.append(ext)

                rename_dict = {mapped_cols['code']: 'code'}
                if mapped_cols['description']: rename_dict[mapped_cols['description']] = 'desc'
                if mapped_cols['quantity']: rename_dict[mapped_cols['quantity']] = 'qty'

                upload_df = raw_df.rename(columns=rename_dict)
                
                if 'qty' not in upload_df.columns:
                    upload_df['qty'] = 1
                else:
                    upload_df['qty'] = pd.to_numeric(upload_df['qty'], errors='coerce').fillna(1).astype(int)

                if 'desc' not in upload_df.columns:
                    upload_df['desc'] = ""

                master_df = self.master_tab.df
                if master_df is None or master_df.empty:
                    QMessageBox.warning(self, "No Master Data", "Please load the master catalogue before bulk uploading offers.")
                    return

                master_codes = set(master_df['code'].astype(str).str.replace('-', '').str.strip().str.lower())
                
                # Pre-build a fast lookup dictionary for descriptions from master_df
                # This is much faster than boolean masking inside the loop
                master_desc_lookup = {}
                for _, m_row in master_df.iterrows():
                    m_code_clean = str(m_row.get('code', '')).replace('-', '').strip().lower()
                    if m_code_clean:
                        master_desc_lookup[m_code_clean] = str(m_row.get('description', ''))

                valid_rows = []
                skipped_codes = []

                for _, row in upload_df.iterrows():
                    code_val = str(row.get('code', '')).strip() if pd.notna(row.get('code')) else ""
                    desc_val = str(row.get('desc', '')).strip() if pd.notna(row.get('desc')) else ""

                    if (not code_val or code_val.lower() == 'nan') and desc_val and desc_val.lower() != 'nan':
                        sec_data = {"code": "", "desc": desc_val, "qty": "", "is_section": True}
                        for ext in self.extra_columns:
                            sec_data[ext] = ""
                        self.offer_data.append(sec_data)
                        valid_rows.append(desc_val)
                        continue

                    if not code_val or code_val.lower() == 'nan':
                        continue

                    code_clean = code_val.replace('-', '').lower()

                    if code_clean not in master_codes:
                        skipped_codes.append(code_val)
                        continue
                        
                    # If description is missing, pull from master file
                    if not desc_val:
                        desc_val = master_desc_lookup.get(code_clean, "")

                    qty_val = int(row['qty']) if pd.notna(row['qty']) else 1
                    row_data = {"code": code_val, "desc": desc_val, "qty": qty_val, "is_section": False}
                    for ext in new_extras:
                        if ext in upload_df.columns and pd.notna(row[ext]):
                            raw_val = row[ext]
                            if isinstance(raw_val, float) and raw_val.is_integer():
                                val = str(int(raw_val)).strip()
                            else:
                                val = str(raw_val).strip()
                        else:
                            val = ""
                        row_data[ext] = val
                    self.offer_data.append(row_data)
                    valid_rows.append(code_val)

                if valid_rows:
                    if hasattr(self, 'master_tab') and hasattr(self.master_tab, 'save_offer_list'):
                        self.master_tab.save_offer_list()
                    self.refresh_offer_table()

                if skipped_codes:
                    QMessageBox.warning(
                        self,
                        "Skipped Invalid Codes",
                        f"The following codes were skipped because they were not found in the master catalogue:\n{', '.join(sorted(set(skipped_codes)))}"
                    )
                else:
                    QMessageBox.information(self, "Bulk Upload Complete", "All uploaded items and sections were added.")
            except Exception as e:
                QMessageBox.critical(self, "Bulk Upload Error", f"Failed to process uploaded file: {e}")

    def get_logical_col_key(self, logical_col):
        if logical_col == 1: return "#", "sr_no"
        elif logical_col == 2: return "CODE", "code"
        elif logical_col == 3: return "DESCRIPTION", "desc"
        elif logical_col == 4: return "QTY", "qty"
        elif logical_col >= 5: return self.extra_columns[logical_col - 5].upper(), self.extra_columns[logical_col - 5]
        return None, None

    def export_excel(self):
        if not self.offer_data:
            return
        path, _ = QFileDialog.getSaveFileName(self, "Save Excel", "Offer_List.xlsx", "Excel Files (*.xlsx)")
        if path:
            with show_loading(self, "Generating Excel..."):
                header = self.offer_table.horizontalHeader()
                visual_cols = [header.logicalIndex(v) for v in range(1, self.offer_table.columnCount())]
                
                export_list = []
                serial_counter = 1
                for data in self.offer_data:
                    export_dict = {}
                    for log_col in visual_cols:
                        header_name, dict_key = self.get_logical_col_key(log_col)
                        if header_name:
                            if data.get("is_section", False):
                                serial_counter = 1
                                export_dict[header_name] = data["desc"] if dict_key == "desc" else ""
                            else:
                                if dict_key == "sr_no":
                                    export_dict[header_name] = str(serial_counter)
                                else:
                                    export_dict[header_name] = data.get(dict_key, "")
                    if not data.get("is_section", False):
                        serial_counter += 1
                    export_list.append(export_dict)

                df = pd.DataFrame(export_list)
                df.to_excel(path, index=False)
                
                try:
                    import openpyxl
                    from openpyxl.styles import PatternFill, Alignment, Font
                    
                    wb = openpyxl.load_workbook(path)
                    ws = wb.active
                    
                    for idx, data in enumerate(self.offer_data):
                        excel_row = idx + 2
                        if data.get("is_section", False):
                            for col_idx in range(1, len(visual_cols) + 1):
                                cell = ws.cell(row=excel_row, column=col_idx)
                                cell.alignment = Alignment(horizontal="center", vertical="center")
                                cell.fill = PatternFill(start_color="D9D9D9", end_color="D9D9D9", fill_type="solid")
                                cell.font = Font(bold=True)
                            
                    wb.save(path)
                except Exception as e:
                    print(f"Failed to style Excel: {e}")

            QMessageBox.information(self, "Success", "Excel Exported Successfully.")

    def _find_image_path(self, code):
        master_df = self.master_tab.df
        if master_df is None or master_df.empty:
            return None
        mask = master_df['code'].astype(str).str.replace('-', '').str.strip().str.lower() == str(code).replace('-', '').strip().lower()
        match = master_df[mask]
        if match.empty:
            return None
        img_path = str(match.iloc[0].get('local_image_path', ''))
        return img_path if img_path and os.path.exists(img_path) else None

    def _scaled_pdf_image(self, img_path, max_w_mm=22, max_h_mm=15):
        max_w, max_h = max_w_mm * mm, max_h_mm * mm
        try:
            img = PILImage.open(img_path)
            img.load()
        except Exception:
            return Paragraph("N/A", ParagraphStyle("na", fontSize=7, alignment=TA_CENTER))

        w, h = img.size
        ratio_upright = min(max_w / w, max_h / h)
        ratio_rotated = min(max_w / h, max_h / w)

        if ratio_rotated > ratio_upright * 1.15:
            img = img.rotate(-90, expand=True)
            w, h = img.size
            ratio = ratio_rotated
        else:
            ratio = ratio_upright

        buf = io.BytesIO()
        img.convert("RGB").save(buf, format="PNG")
        buf.seek(0)
        return RLImage(buf, width=w * ratio, height=h * ratio)

    def export_pdf(self):
        if not self.offer_data:
            return

        settings_dialog = PdfSettingsDialog(self)
        if settings_dialog.exec() != QDialog.DialogCode.Accepted:
            return
        settings = settings_dialog.get_settings()

        path, _ = QFileDialog.getSaveFileName(self, "Save PDF", "Offer_List.pdf", "PDF Files (*.pdf)")
        if not path:
            return

        with show_loading(self, "Generating PDF, this may take a moment..."):
            try:
                try:
                    pdfmetrics.registerFont(TTFont('Arial', 'C:\\Windows\\Fonts\\arial.ttf'))
                    pdfmetrics.registerFont(TTFont('Arial-Bold', 'C:\\Windows\\Fonts\\arialbd.ttf'))
                    f_norm = 'Arial'
                    f_bold = 'Arial-Bold'
                except Exception:
                    f_norm = 'Helvetica'
                    f_bold = 'Helvetica-Bold'

                styles = getSampleStyleSheet()
                cell_style = ParagraphStyle("cell", parent=styles["Normal"], fontSize=8, leading=10, fontName=f_norm)
                qty_style = ParagraphStyle("qty", parent=cell_style, alignment=TA_CENTER)
                header_cell_style = ParagraphStyle(
                    "hcell", parent=styles["Normal"], fontSize=8, leading=10,
                    textColor=colors.white, fontName=f_bold)

                doc = SimpleDocTemplate(
                    path, pagesize=A4,
                    leftMargin=15 * mm, rightMargin=15 * mm,
                    topMargin=15 * mm, bottomMargin=15 * mm,
                    title=settings.get("reference", "Offer")
                )

                story = []
                if settings.get('header'):
                    story.append(Paragraph(shape_arabic(settings['header']), ParagraphStyle(
                        "h1", parent=styles["Heading1"], fontSize=16, spaceAfter=2, fontName=f_bold)))
                if settings.get('subheader'):
                    story.append(Paragraph(shape_arabic(settings['subheader']), ParagraphStyle(
                        "h2", parent=styles["Heading2"], fontSize=12, textColor=colors.grey, spaceAfter=2, fontName=f_norm)))
                if settings.get('reference'):
                    story.append(Paragraph(shape_arabic(settings['reference']), ParagraphStyle(
                        "h3", parent=styles["Normal"], fontSize=10, spaceAfter=10, fontName=f_norm)))
                story.append(Spacer(1, 6))

                header = self.offer_table.horizontalHeader()
                visual_cols = [header.logicalIndex(v) for v in range(1, self.offer_table.columnCount())]
                
                header_row = []
                col_widths = []
                
                usable_width_mm = 180
                total_text_cols_width = 0
                
                text_col_widths_mm = []
                for log_col in visual_cols:
                    header_name, dict_key = self.get_logical_col_key(log_col)
                    if not header_name:
                        text_col_widths_mm.append(0)
                        continue
                    
                    shaped_header = shape_arabic(header_name)
                    max_w = stringWidth(shaped_header, f_bold, 8) / mm
                    
                    for data in self.offer_data:
                        if not data.get("is_section", False):
                            text = str(data.get(dict_key, ""))
                            shaped_text = shape_arabic(text)
                            w = stringWidth(shaped_text, f_norm, 8) / mm
                            if w > max_w:
                                max_w = w
                                
                    col_w = max_w + 6
                    if header_name == "DESCRIPTION":
                        col_w = min(col_w, 85.0)
                    elif header_name == "CODE":
                        col_w = min(col_w, 45.0)
                    else:
                        col_w = min(col_w, 35.0)

                    text_col_widths_mm.append(col_w)
                    total_text_cols_width += col_w

                image_col_width = max(25.0, usable_width_mm - total_text_cols_width)
                
                for i, log_col in enumerate(visual_cols):
                    header_name, _ = self.get_logical_col_key(log_col)
                    if header_name:
                        header_row.append(Paragraph(shape_arabic(header_name), header_cell_style))
                        col_widths.append(text_col_widths_mm[i] * mm)
                
                header_row.append(Paragraph("Image", header_cell_style))
                col_widths.append(image_col_width * mm)
                
                table_data = [header_row]
                
                dynamic_styles = []

                serial_counter = 1
                for row_idx, data in enumerate(self.offer_data, start=1):
                    if data.get("is_section", False):
                        serial_counter = 1
                        sec_p = Paragraph(f"<b>{shape_arabic(data['desc'])}</b>", ParagraphStyle("sec", parent=cell_style, alignment=TA_CENTER, fontName=f_bold))
                        row_data = [sec_p] + [""] * len(visual_cols)
                        table_data.append(row_data)
                        dynamic_styles.append(("SPAN", (0, row_idx), (-1, row_idx)))
                        dynamic_styles.append(("BACKGROUND", (0, row_idx), (-1, row_idx), colors.HexColor("#DDE3EC")))
                        dynamic_styles.append(("ALIGN", (0, row_idx), (-1, row_idx), "CENTER"))
                    else:
                        img_cell = ""
                        if settings['include_images']:
                            img_path = self._find_image_path(data['code'])
                            avail_img_w = max(5, image_col_width - 4)
                            img_cell = self._scaled_pdf_image(img_path, max_w_mm=avail_img_w, max_h_mm=25) if img_path else Paragraph("—", cell_style)

                        row_data = []
                        for log_col in visual_cols:
                            header_name, dict_key = self.get_logical_col_key(log_col)
                            if header_name:
                                if dict_key == "sr_no":
                                    row_data.append(Paragraph(str(serial_counter), qty_style))
                                else:
                                    style = qty_style if header_name == "QTY" else cell_style
                                    row_data.append(Paragraph(shape_arabic(str(data.get(dict_key, ""))), style))
                        
                        row_data.append(img_cell)
                        table_data.append(row_data)
                        serial_counter += 1

                tbl = Table(table_data, colWidths=col_widths, repeatRows=1)
                
                base_style = [
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2C3E50")),
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#CCCCCC")),
                    ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                    ("TOPPADDING", (0, 0), (-1, -1), 4),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                    ("LEFTPADDING", (0, 0), (-1, -1), 5),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 5),
                ]
                
                tbl.setStyle(TableStyle(base_style + dynamic_styles))
                story.append(tbl)

                doc.build(story)

            except Exception as e:
                QMessageBox.critical(self, "Error", f"Generation failed: {str(e)}")
                return
                
        QMessageBox.information(self, "Success", "PDF Generated Successfully.")

    def export_pptx(self):
        if not self.offer_data:
            return

        import os
        import sys
        
        # Check if running as a PyInstaller bundled executable
        if getattr(sys, 'frozen', False):
            base_path = sys._MEIPASS
        else:
            base_path = os.path.abspath(".")
            
        template_path = os.path.join(base_path, "template.pptx")
        
        if not os.path.exists(template_path):
            template_path, _ = QFileDialog.getOpenFileName(self, "Select PPTX Template", "", "PowerPoint Files (*.pptx)")
            if not template_path:
                return

        save_path, _ = QFileDialog.getSaveFileName(self, "Save PPTX", "Offer_Presentation.pptx", "PowerPoint Files (*.pptx)")
        if not save_path:
            return

        with show_loading(self, "Generating Presentation..."):
            try:
                from pptx import Presentation
                from pptx.enum.shapes import PP_PLACEHOLDER
                import re
                from core.constants import BROCHURE_HIERARCHY

                prs = Presentation(template_path)
                
                if len(prs.slides) == 0:
                    raise ValueError("Template must contain at least one slide.")
                    
                dummy_slide = prs.slides[-1]
                dummy_rId = prs.slides._sldIdLst[-1].rId
                dummy_sld_elem = prs.slides._sldIdLst[-1]
                slide_layout = dummy_slide.slide_layout
                
                master_df = self.master_tab.df

                def get_disciplines(brochure_str):
                    if not brochure_str or pd.isna(brochure_str): return ""
                    found_disciplines = set()
                    brocs = [b.strip() for b in str(brochure_str).replace(';', ',').split(',')]
                    for b in brocs:
                        for main_cat, sub_list in BROCHURE_HIERARCHY.items():
                            if b in sub_list or b == main_cat:
                                found_disciplines.add(main_cat)
                    return ", ".join(sorted(found_disciplines)) if found_disciplines else "General"

                def get_sizes(desc):
                    if not desc: return ""
                    matches = re.findall(r'\b\d+(?:\.\d+)?\s*(?:mm|cm|inch|m)\b', desc, re.IGNORECASE)
                    return ", ".join(matches) if matches else "Standard"
                
                def apply_context(text, context):
                    for key, val in context.items():
                        text = text.replace(f"{{{{{key}}}}}", str(val))
                    return text

                # First, enrich offer data
                items_to_process = []
                
                for data in self.offer_data:
                    if data.get("is_section", False):
                        continue
                        
                    code = data.get("code", "")
                    desc = data.get("desc", "")
                    brochures = ""
                    disciplines = ""
                    sizes = ""
                    
                    if master_df is not None and not master_df.empty:
                        mask = master_df['code'].astype(str).str.replace('-', '').str.replace(' ', '').str.lower() == str(code).replace('-', '').replace(' ', '').lower()
                        match = master_df[mask]
                        if not match.empty:
                            master_row = match.iloc[0]
                            desc = master_row.get("description", desc)
                            brochures_val = master_row.get("brochures", "")
                            brochures = brochures_val if pd.notna(brochures_val) else ""
                            disciplines = get_disciplines(brochures)
                            sizes = get_sizes(desc)
                            
                        # We remove family logic for now
                                
                    item = {
                        "code": code, "desc": desc, "brochures": brochures,
                        "disciplines": disciplines, "sizes": sizes
                    }
                    items_to_process.append(item)

                for item in items_to_process:
                    context = {
                        "CODE": item["code"],
                        "DESCRIPTION": item["desc"],
                        "BROCHURES": item["brochures"],
                        "DISCIPLINES": item["disciplines"],
                        "SIZES": item["sizes"],
                        "sizes": item["sizes"],
                        "VARIATIONS": ""
                    }
                    
                    img_path = self._find_image_path(item["code"])
                    
                    # Convert comma/semicolon-separated lists into bullet points
                    for k in ["BROCHURES", "DISCIPLINES", "SIZES", "sizes"]:
                        if context[k]:
                            parts = [x.strip() for x in str(context[k]).replace(';', ',').split(',') if x.strip()]
                            if len(parts) > 1:
                                context[k] = "\n".join(f"• {p}" for p in parts)
                    
                    slide = prs.slides.add_slide(slide_layout)
                    
                    import copy
                    from PIL import Image
                    
                    # Copy background if it exists on the dummy slide explicitly
                    try:
                        dummy_bg = dummy_slide.element.xpath('./p:cSld/p:bg')
                        if dummy_bg:
                            existing_bg = slide.element.xpath('./p:cSld/p:bg')
                            if existing_bg:
                                slide.element.xpath('./p:cSld')[0].remove(existing_bg[0])
                            
                            new_bg = copy.deepcopy(dummy_bg[0])
                            
                            # Map relationships for background if it has a picture fill
                            for blip in new_bg.xpath('.//a:blip'):
                                rId = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                                if rId:
                                    try:
                                        rel = dummy_slide.part.rels[rId]
                                        target_part = rel.target_part
                                        new_rId = slide.part.relate_to(target_part, rel.reltype)
                                        blip.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed', new_rId)
                                    except KeyError:
                                        pass
                            
                            slide.element.xpath('./p:cSld')[0].insert(0, new_bg)
                    except Exception:
                        pass
                    
                    img_box = None
                    for shape in dummy_slide.shapes:
                        if shape.is_placeholder:
                            continue # Placeholders are handled by the layout
                            
                        # Check for {{IMAGE}} tag
                        is_img_tag = False
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if "{{IMAGE}}" in run.text:
                                        img_box = (shape.left, shape.top, shape.width, shape.height)
                                        is_img_tag = True
                                        break
                                if is_img_tag: break
                                
                        if is_img_tag:
                            continue # Do not copy the {{IMAGE}} text box to the new slide
                            
                        el = shape.element
                        new_el = copy.deepcopy(el)
                        
                        # Fix relationships for any images in the copied shape (including inside groups)
                        for blip in new_el.xpath('.//a:blip'):
                            rId = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                            if rId:
                                try:
                                    rel = dummy_slide.part.rels[rId]
                                    target_part = rel.target_part
                                    new_rId = slide.part.relate_to(target_part, rel.reltype)
                                    blip.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed', new_rId)
                                except KeyError:
                                    pass
                                    
                        slide.shapes._spTree.append(new_el)
                    
                    # 2. Process all shapes on the new slide
                    for shape in slide.shapes:
                        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                            # Keep fallback for Slide Master placeholders
                            if img_path and not img_box:
                                try:
                                    shape.insert_picture(img_path)
                                except Exception:
                                    pass
                        elif shape.has_text_frame:
                            # Replace text in the shape
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.text = apply_context(run.text, context)
                                    
                    # 3. Insert Custom {{IMAGE}} bounding box
                    if img_path and img_box:
                        try:
                            left, top, width, height = img_box
                            with Image.open(img_path) as img:
                                orig_w, orig_h = img.size
                            
                            aspect_orig = orig_w / orig_h
                            aspect_box = width / height
                            
                            if aspect_orig > aspect_box:
                                # Image is wider than box, constrain to width
                                new_w = width
                                new_h = width / aspect_orig
                            else:
                                # Image is taller than box, constrain to height
                                new_h = height
                                new_w = height * aspect_orig
                                
                            # Center it in the box
                            new_left = left + (width - new_w) / 2
                            new_top = top + (height - new_h) / 2
                            
                            slide.shapes.add_picture(img_path, new_left, new_top, new_w, new_h)
                        except Exception:
                            pass

                try:
                    prs.part.drop_rel(dummy_rId)
                    prs.slides._sldIdLst.remove(dummy_sld_elem)
                except Exception:
                    pass

                prs.save(save_path)
                
            except Exception as e:
                QMessageBox.critical(self, "Error", f"PPTX Generation failed: {str(e)}")
                return
                
        # Outside show_loading
        msgBox = QMessageBox(self)
        msgBox.setWindowTitle("Success")
        msgBox.setText("PowerPoint Presentation Generated Successfully.")
        msgBox.setInformativeText("Would you like to save a PDF copy as well, or just open the PPTX?")
        
        pdf_btn = msgBox.addButton("Convert to PDF", QMessageBox.ButtonRole.ActionRole)
        open_btn = msgBox.addButton("Open PPTX", QMessageBox.ButtonRole.ActionRole)
        close_btn = msgBox.addButton("Close", QMessageBox.ButtonRole.RejectRole)
        
        msgBox.exec()
        
        if msgBox.clickedButton() == pdf_btn:
            pdf_path = save_path.rsplit('.', 1)[0] + ".pdf"
            with show_loading(self, "Converting PPTX to PDF..."):
                try:
                    import comtypes.client
                    comtypes.CoInitialize()
                    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
                    deck = powerpoint.Presentations.Open(os.path.abspath(save_path), WithWindow=False)
                    deck.SaveAs(os.path.abspath(pdf_path), 32)
                    deck.Close()
                    powerpoint.Quit()
                except Exception as ex:
                    try: powerpoint.Quit() 
                    except: pass
                    QMessageBox.critical(self, "Error", f"Failed to convert to PDF:\n{str(ex)}")
                    return
            
            QMessageBox.information(self, "Success", f"PDF saved successfully at:\n{pdf_path}")
            try:
                os.startfile(os.path.abspath(pdf_path))
            except Exception:
                pass
                
        elif msgBox.clickedButton() == open_btn:
            try:
                os.startfile(os.path.abspath(save_path))
            except Exception:
                pass
